import os
from pathlib import Path
from typing import List
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from utils.RAGLogger import RAGLogger

import time

class DocumentLoader:
    def __init__(self, input_dir: str, output_dir: str):
        self.logger = RAGLogger('DocumentLoader')
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
       

    def load_documents(self) -> List[Path]:
        """Recursively load documents from the input directory and subfolders."""
        supported_exts = {'.pdf', '.docx', '.pptx', '.xlsx'}
        documents = []
        all_files = list(self.input_dir.rglob('*'))
        self.logger.info(
            f"Scanning directory for documents",
            total_files_found=len(all_files),
            supported_extensions=list(supported_exts),
            scan_directory=str(self.input_dir)
        )
        for file in all_files:
            self.logger.debug(
                f"Examining file: {file.name}",
                file_path = str(file),
                file_suffix=file.suffix,
                is_supported=file.suffix.lower() in supported_exts
            )
            if file.is_file() and file.suffix.lower() in supported_exts:
                documents.append(file)

        self.logger.info(
            f"Document discovery completed",
            supported_documents_found=len(documents),
            file_types_found=list(set(doc.suffix for doc in documents))
        )
        return documents

    def parse_pdf(self, file_path: Path) -> str:
        """Extract text from a PDF file."""
        try:
            self.logger.debug(f"Starting PDF parsing: {file_path.name}")
            reader = PdfReader(file_path)

            # Log PDF details
            self.logger.debug(
                f"PDF metadata extracted",
                file_name=file_path.name,
                total_pages=len(reader.pages)
            )

            text = "\n".join(page.extract_text() for page in reader.pages)

            self.logger.debug(
                f"PDF text extraction completed",
                file_name=file_path.name,
                extracted_characters=len(text)
            )
            return text
        
        except Exception as e:
            self.logger.error(
                f"PDF parsing failed: {file_path.name}",
                file_path=str(file_path),
                error_type=type(e).__name__,
                error_message=str(e)
            )
            raise

        

    def parse_docx(self, file_path: Path) -> str:
        """Extract text from a Word document."""
        try:
            self.logger.debug(f"Starting DOCX parsing: {file_path.name}")
            doc = Document(file_path)
            
            # Log DOCX details
            self.logger.debug(
                f"DOCX metadata extracted",
                file_name=file_path.name,
                total_paragraphs=len(doc.paragraphs)
            )
            
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            
            self.logger.debug(
                f"DOCX text extraction completed",
                file_name=file_path.name,
                extracted_characters=len(text)
            )
            return text
            
        except Exception as e:
            self.logger.error(
                f"DOCX parsing failed: {file_path.name}",
                file_path=str(file_path),
                error_type=type(e).__name__,
                error_message=str(e)
            )
            raise

    def parse_pptx(self, file_path: Path) -> str:
        """Extract text from a PowerPoint file."""
        try:
            self.logger.debug(f"Starting PPTX parsing: {file_path.name}")
            presentation = Presentation(file_path)
            
            # Log PPTX details
            self.logger.debug(
                f"PPTX metadata extracted",
                file_name=file_path.name,
                total_slides=len(presentation.slides)
            )
            
            text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, 'text'))
            
            self.logger.debug(
                f"PPTX text extraction completed",
                file_name=file_path.name,
                extracted_characters=len(text)
            )
            return text
            
        except Exception as e:
            self.logger.error(
                f"PPTX parsing failed: {file_path.name}",
                file_path=str(file_path),
                error_type=type(e).__name__,
                error_message=str(e)
            )
            raise

    def parse_xlsx(self, file_path: Path) -> str:
        """Extract text from an Excel file."""
        try:
            self.logger.debug(f"Starting XLSX parsing: {file_path.name}")
            workbook = load_workbook(file_path, data_only=True)
            
            # Log XLSX details
            self.logger.debug(
                f"XLSX metadata extracted",
                file_name=file_path.name,
                total_worksheets=len(workbook.worksheets),
                worksheet_names=[sheet.title for sheet in workbook.worksheets]
            )
            
            text = "\n".join(
                "\n".join(str(cell.value) for cell in row if cell.value is not None)
                for sheet in workbook.worksheets
                for row in sheet.iter_rows()
            )
            
            self.logger.debug(
                f"XLSX text extraction completed",
                file_name=file_path.name,
                extracted_characters=len(text)
            )
            return text
            
        except Exception as e:
            self.logger.error(
                f"XLSX parsing failed: {file_path.name}",
                file_path=str(file_path),
                error_type=type(e).__name__,
                error_message=str(e)
            )
            raise

    def save_text(self, file_name: str, text: str):
        """Save extracted text to the output directory."""
        try:

            # Replace path separators with underscores for unique file names
            safe_name = file_name.replace(os.sep, '_')
            output_path = self.output_dir / f"{safe_name}.txt"

            self.logger.debug(
                f"Saving extracted text to file",
                original_filename = file_name,
                safe_filename=safe_name,
                output_path=str(output_path),
                text_length=len(text)
            )

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)

            # Verify file was created

            if output_path.exists():
                file_size = output_path.stat().st_size
                self.logger.debug(
                    f"Text file saved successfully",
                    output_file=str(output_path),
                    file_size_bytes=file_size,
                    text_length=len(text)
                )
            else:
                self.logger.error(f"File save verification failed: {output_path}")
        
        except Exception as e:
            self.logger.error(
                f"Failed to save text file",
                filename=file_name,
                error_message=str(e),
                text_length=len(text)
            )
            raise

    def process_documents(self) -> dict:
        """
        Process all documents and return results summary.
        Continues processing even if individual files fail.
        """
        start_time = time.time()
        files = self.load_documents()

        self.logger.info(
            f"Starting document processing batch",
            total_files=len(files),
            file_types=list(set(f.suffix for f in files)),
            batch_id=f"batch_{int(start_time)}"
        )

        results = {
            "successful_files": [],
            "failed_files": [],
            "total_files": len(files),
            "processing_time": 0
        }

        for file in files:
            file_start_time = time.time()
            
            try:
                # Parse the file using the new method
                text, metadata = self.parse_single_file(file)
                processing_time = time.time() - file_start_time
                
                if metadata["status"] == "success":
                    # Save the text file
                    rel_path = str(file.relative_to(self.input_dir)).replace(os.sep, '_')
                    self.save_text(rel_path, text)
                    
                    results["successful_files"].append({
                        "file_path": file,
                        "relative_path": rel_path,
                        "text": text,
                        "metadata": metadata,
                        "processing_time": processing_time
                    })
                    
                    self.logger.info(
                        f"Successfully processed: {file.name}",
                        processing_time_seconds=round(processing_time, 3),
                        text_length=len(text)
                    )
                else:
                    results["failed_files"].append({
                        "file_path": file,
                        "metadata": metadata,
                        "processing_time": processing_time
                    })
                    
            except Exception as e:
                # This catches any unexpected errors not handled in parse_single_file
                processing_time = time.time() - file_start_time
                self.logger.error(f"Unexpected error processing {file.name}: {e}")
                
                results["failed_files"].append({
                    "file_path": file,
                    "metadata": {
                        "status": "failed",
                        "error_message": f"Unexpected error: {str(e)}",
                        "text_length": 0
                    },
                    "processing_time": processing_time
                })

        total_time = time.time() - start_time
        results["processing_time"] = total_time
        
        self.logger.info(
            f"Document processing batch completed",
            total_files=results["total_files"],
            successful_files=len(results["successful_files"]),
            failed_files=len(results["failed_files"]),
            success_rate=round((len(results["successful_files"]) / results["total_files"]) * 100, 2) if results["total_files"] > 0 else 0,
            total_processing_time_seconds=round(total_time, 3)
        )

        return results
    
    def get_file_metadata(self, file_path: Path) -> dict:
        """Get file metadata for manifest tracking"""
        try:
            stat = file_path.stat()
            return {
                "file_size": stat.st_size,
                "last_modified": stat.st_mtime,
                "file_type": file_path.suffix.lower()
            }
        except Exception as e:
            self.logger.error(f"Error getting metadata for {file_path}: {e}")
            return {"file_size": 0, "last_modified": 0, "file_type": "unknown"}

    def parse_single_file(self, file_path: Path) -> tuple[str, dict]:
        """
        Parse a single file and return (text, metadata).
        Returns empty string and error info if parsing fails.
        """
        try:
            self.logger.info(f"Parsing file: {file_path.name}")
            
            if file_path.suffix.lower() == '.pdf':
                text = self.parse_pdf(file_path)
            elif file_path.suffix.lower() == '.docx':
                text = self.parse_docx(file_path)
            elif file_path.suffix.lower() == '.pptx':
                text = self.parse_pptx(file_path)
            elif file_path.suffix.lower() == '.xlsx':
                text = self.parse_xlsx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_path.suffix}")
            
            metadata = self.get_file_metadata(file_path)
            metadata.update({
                "status": "success",
                "error_message": None,
                "text_length": len(text)
            })
            
            return text, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to parse {file_path.name}: {e}")
            metadata = self.get_file_metadata(file_path)
            metadata.update({
                "status": "failed",
                "error_message": str(e),
                "text_length": 0
            })
            return "", metadata


if __name__ == "__main__":
    import os
    # Set project root to three levels up from this script (sharepoint-rag-bot)
    root_dir = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
    input_dir = os.path.join(root_dir, "data", "raw")
    output_dir = input_dir
    print(f"Using input_dir: {input_dir}")
    loader = DocumentLoader(input_dir, output_dir)
    loader.process_documents()
